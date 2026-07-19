#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成"课堂实录标注与分析系统"项目成果汇报PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ========== 颜色方案 ==========
PRIMARY = RGBColor(0x4F, 0x46, 0xE5)    # 靛蓝
ACCENT = RGBColor(0x06, 0xB6, 0xD4)      # 青色
DARK = RGBColor(0x1E, 0x29, 0x3B)        # 深色文字
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)       # 浅色背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color):
    """设置幻灯片纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    """添加色块形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_card(slide, left, top, width, height, title, content, title_color=PRIMARY):
    """添加卡片"""
    card = add_shape(slide, left, top, width, height, WHITE)
    card.shadow.inherit = False
    # 标题
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2),
                 width - Inches(0.6), Inches(0.5), title,
                 font_size=16, color=title_color, bold=True)
    # 内容
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7),
                 width - Inches(0.6), height - Inches(0.9), content,
                 font_size=13, color=DARK)


def add_top_bar(slide, title_text):
    """添加顶部标题栏"""
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
                 title_text, font_size=30, color=WHITE, bold=True)


def add_page_number(slide, num, total):
    """添加页码"""
    add_text_box(slide, Inches(12), Inches(7.05), Inches(1.2), Inches(0.4),
                 f'{num}/{total}', font_size=11, color=GRAY, alignment=PP_ALIGN.RIGHT)


TOTAL = 12

# =====================================================
# 第1页：封面
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), PRIMARY)
add_shape(slide, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), PRIMARY)
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             '课堂实录标注与分析系统', font_size=42, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.1), Inches(11), Inches(0.7),
             '—— 基于 Whisper 语音识别的 AI 智能教学行为分析平台',
             font_size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)
# 分隔线
add_shape(slide, Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.03), ACCENT)
add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.6),
             '项目实践成果汇报', font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             '2026年6月', font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 1, TOTAL)

# =====================================================
# 第2页：项目背景与目标
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT)
add_top_bar(slide, '项目背景与目标')

# 背景
add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.5),
         '📋 项目背景', 
         '传统课堂评课依赖人工观察记录，耗时费力且主观性强。\n'
         '随着 AI 技术的发展，利用语音识别和自然语言处理技术\n'
         '自动分析课堂教学行为成为可能，能够为教师提供客观、\n'
         '量化的教学反馈，助力教学质量的提升。')

add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.5),
         '🎯 项目目标',
         '构建一个 Web 平台，支持上传课堂实录视频，自动进行\n'
         '语音识别和教学行为标注，生成多维度的数据分析报告，\n'
         '帮助教师全面了解课堂教学结构、师生互动模式和教学\n'
         '节奏特点，为教学反思和改进提供数据支撑。')

# 核心价值
add_shape(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3), WHITE)
add_text_box(slide, Inches(0.8), Inches(4.7), Inches(2), Inches(0.5),
             '💡 核心价值', font_size=18, color=PRIMARY, bold=True)

values = [
    ('🎙️', '自动化', 'AI自动完成语音识别\n和行为标注'),
    ('📊', '多维度', '弗兰德斯、S-T等多维\n教学分析模型'),
    ('⏱️', '高效率', '数分钟完成人工需要\n数小时的分析工作'),
    ('📈', '可量化', '将主观评课转化为\n客观数据指标'),
]
for i, (icon, title, desc) in enumerate(values):
    x = Inches(1 + i * 3)
    add_text_box(slide, x, Inches(5.3), Inches(2.8), Inches(0.4),
                 f'{icon} {title}', font_size=15, color=PRIMARY, bold=True)
    add_text_box(slide, x, Inches(5.8), Inches(2.8), Inches(0.8),
                 desc, font_size=12, color=DARK)

add_page_number(slide, 2, TOTAL)

# =====================================================
# 第3页：技术架构
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT)
add_top_bar(slide, '技术架构总览')

# 前端
add_shape(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(2.2), WHITE)
add_shape(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(0.5), PRIMARY)
add_text_box(slide, Inches(0.5), Inches(1.65), Inches(3.8), Inches(0.5),
             '🖥️  前端', font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(3.3), Inches(1.3),
             '• Vite 构建工具\n'
             '• 原生 JavaScript（ES Module）\n'
             '• Chart.js 数据可视化\n'
             '• Canvas 时间轴绘制\n'
             '• 单页应用（SPA）架构',
             font_size=13, color=DARK)

# 后端
add_shape(slide, Inches(4.8), Inches(1.6), Inches(3.8), Inches(2.2), WHITE)
add_shape(slide, Inches(4.8), Inches(1.6), Inches(3.8), Inches(0.5), ACCENT)
add_text_box(slide, Inches(4.8), Inches(1.65), Inches(3.8), Inches(0.5),
             '⚙️  后端', font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.1), Inches(2.3), Inches(3.3), Inches(1.3),
             '• Node.js + Express\n'
             '• Multer 文件上传\n'
             '• FFmpeg 音视频处理\n'
             '• Whisper.cpp 语音识别\n'
             '• 关键词匹配 + 片段合并算法',
             font_size=13, color=DARK)

# AI层
add_shape(slide, Inches(9.1), Inches(1.6), Inches(3.8), Inches(2.2), WHITE)
add_shape(slide, Inches(9.1), Inches(1.6), Inches(3.8), Inches(0.5), ORANGE)
add_text_box(slide, Inches(9.1), Inches(1.65), Inches(3.8), Inches(0.5),
             '🤖  AI 引擎', font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.4), Inches(2.3), Inches(3.3), Inches(1.3),
             '• OpenAI Whisper 模型\n'
             '• 中文语音转文本\n'
             '• 教学行为关键词库\n'
             '• 智能片段合并策略\n'
             '• 多维度数据分析算法',
             font_size=13, color=DARK)

# 箭头
add_text_box(slide, Inches(4.3), Inches(2.4), Inches(0.5), Inches(0.5),
             '→', font_size=30, color=PRIMARY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.6), Inches(2.4), Inches(0.5), Inches(0.5),
             '→', font_size=30, color=ACCENT, alignment=PP_ALIGN.CENTER)

# 底部：数据流
add_shape(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.8), WHITE)
add_text_box(slide, Inches(0.8), Inches(4.4), Inches(2), Inches(0.5),
             '📡 数据流转', font_size=18, color=PRIMARY, bold=True)

flow_steps = [
    ('上传视频', '用户上传\n课堂实录'),
    ('提取音频', 'FFmpeg 提取\nWAV 音频'),
    ('语音识别', 'Whisper 转文本\n带时间戳'),
    ('AI 标注', '关键词匹配\n生成教学标签'),
    ('数据分析', '多维模型\n聚合计算'),
    ('结果展示', '图表报告\n可视化呈现'),
]
for i, (title, desc) in enumerate(flow_steps):
    x = Inches(0.8 + i * 2.05)
    add_shape(slide, x, Inches(5.0), Inches(1.85), Inches(0.4), PRIMARY)
    add_text_box(slide, x, Inches(5.0), Inches(1.85), Inches(0.4),
                 title, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(5.5), Inches(1.85), Inches(0.8),
                 desc, font_size=11, color=DARK, alignment=PP_ALIGN.CENTER)
    if i < 5:
        add_text_box(slide, Inches(x.inches + 1.85), Inches(5.05), Inches(0.2), Inches(0.35),
                     '▸', font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(6.4), Inches(12), Inches(0.5),
             '💡 用户上传视频后，后端自动完成流水线处理，前端实时轮询进度，处理完成后自动加载标注和分析结果',
             font_size=12, color=GRAY)

add_page_number(slide, 3, TOTAL)

# =====================================================
# 第4页：核心功能 - 视频上传与处理
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT)
add_top_bar(slide, '核心功能一：视频上传与智能处理')

add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.8),
         '📤 上传与处理流程',
         '① 支持拖拽或点击上传视频文件（MP4/AVI/MOV）\n'
         '② 前端显示文件名、大小等预览信息\n'
         '③ 后端接收后启动自动化流水线处理\n'
         '④ 前端轮询展示实时处理进度\n'
         '⑤ 处理完成后自动加载标注和分析结果\n\n'
         '处理步骤：视频转码 → 音频提取 → 语音识别 → 智能标注')

add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.8),
         '🔧 关键技术点',
         '• Multer 中间件处理大文件上传\n'
         '• FFmpeg 提取 16kHz 单声道 WAV 音频\n'
         '• FFmpeg 转码 H.264 MP4 适配浏览器播放\n'
         '• 内存 Map 管理任务状态（pending→processing→completed）\n'
         '• 前端 1.5 秒轮询 + 连续错误计数容错\n'
         '• Loading 遮罩实时显示进度文本')

add_card(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
         '🎯 处理流水线架构',
         '┌─────────┐     ┌─────────┐     ┌─────────┐     ┌─────────┐\n'
         '│ 上传视频  │ ──→ │ 提取音频  │ ──→ │ 语音识别  │ ──→ │ AI 标注  │\n'
         '│ (原始MP4) │     │ (WAV)    │     │ (Whisper) │     │ (Tags)   │\n'
         '└─────────┘     └─────────┘     └─────────┘     └─────────┘\n'
         '                                                     ↓\n'
         '                                              ┌─────────┐\n'
         '                                              │ 视频转码  │\n'
         '                                              │ (H.264)  │\n'
         '                                              └─────────┘')

add_page_number(slide, 4, TOTAL)

# =====================================================
# 第5页：核心功能 - AI 智能标注
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT)
add_top_bar(slide, '核心功能二：AI 智能教学行为标注')

add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.8),
         '🧠 标注算法原理',
         '① Whisper 语音识别 → 带时间戳的文本片段\n'
         '② 遍历每个片段进行关键词匹配\n'
         '③ 匹配成功：按规则生成标注标签\n'
         '④ 未匹配片段：合并相邻短片段（≤3秒间隙）\n'
         '⑤ 合并后按时长分配类型（≥20s讲解，<8s丢弃）\n\n'
         '避免碎片化标注，保证标注质量')

add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.8),
         '🏷️ 教学行为分类体系',
         '📘 教师行为：讲解、提问、反馈、示范、总结\n'
         '📗 学生行为：回答、练习、展示、提问\n'
         '📙 互动行为：师生问答、小组合作\n'
         '📕 教学环节：上课、下课、课堂管理\n\n'
         '共建立 12 个关键词分类、80+ 匹配模式\n'
         '涵盖真实课堂的主要教学场景')

add_card(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
         '📋 关键词匹配示例',
         '┌──────────────────┬──────────────────────────────────────┐\n'
         '│ 行为类型          │ 匹配关键词                            │\n'
         '├──────────────────┼──────────────────────────────────────┤\n'
         '│ 教师提问          │ 谁来回答、为什么、怎么、你觉得...         │\n'
         '│ 教师反馈          │ 很好、不错、正确、再想想、还有吗...       │\n'
         '│ 学生练习          │ 练习、试试、做题、练一练...              │\n'
         '│ 小组合作          │ 讨论、小组、交流、同桌、一起...           │\n'
         '│ 课程讲解          │ 上课、下课、安静、看黑板、注意听...       │\n'
         '│ 学生提问          │ 举手、提问、不懂、没听懂、请问...         │\n'
         '└──────────────────┴──────────────────────────────────────┘')

add_page_number(slide, 5, TOTAL)

# =====================================================
# 第6页：核心功能 - 标注编辑与时间轴
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT)
add_top_bar(slide, '核心功能三：标注编辑与时间轴可视化')

add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.8),
         '✏️ 标注管理功能',
         '• 自动标注结果可手动编辑、删除\n'
         '• 支持手动添加新的教学行为标注\n'
         '• 每条标注显示：类型标签、时间段、备注\n'
         '• 不同行为类型用不同颜色区分\n'
         '  教师行为(蓝) / 学生行为(绿) / 互动(橙)\n'
         '• 点击标注自动跳转视频到对应时间点\n'
         '• 标注数据实时同步到分析和报告模块')

add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.8),
         '⏱️ 时间轴可视化',
         '• Canvas 绘制课堂教学时间轴\n'
         '• 按行为类别分层渲染彩色色块\n'
         '• 色块宽度对应该行为的持续时间\n'
         '• 支持鼠标拖拽滚动和滚轮缩放\n'
         '• 悬停色块显示标注详情 Tooltip\n'
         '• 时间轴与视频播放器双向同步\n'
         '• 直观展示整节课的教学节奏')

add_card(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
         '🔄 数据同步机制',
         '┌──────────┐    setTags()    ┌──────────┐    onChange()    ┌──────────┐\n'
         '│ 上传模块  │ ──────────────→ │ DataStore │ ───────────────→ │ 标注模块  │\n'
         '│ (AI标注)  │                 │  (状态仓库) │                 │ (列表渲染) │\n'
         '└──────────┘                 └──────────┘                 └──────────┘\n'
         '                                   │                              \n'
         '                    ┌──────────────┼──────────────┐              \n'
         '                    ↓              ↓              ↓              \n'
         '              ┌──────────┐ ┌──────────┐ ┌──────────┐            \n'
         '              │ 时间轴    │ │ 分析图表  │ │ 报告模块  │            \n'
         '              └──────────┘ └──────────┘ └──────────┘            ')

add_page_number(slide, 6, TOTAL)

# =====================================================
# 第7页：核心功能 - 多维度数据分析
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT)
add_top_bar(slide, '核心功能四：多维度教学数据分析')

add_card(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(2.5),
         '📊 弗兰德斯互动分析',
         '统计教师言语、学生言语、\n'
         '沉默/混乱三类行为的占比，\n'
         '用环形图直观展示课堂互动\n'
         '结构，判断课堂是以教师为\n'
         '中心还是以学生为中心。')

add_card(slide, Inches(4.8), Inches(1.6), Inches(3.8), Inches(2.5),
         '📈 S-T 教学分析',
         '计算教师行为占有率 Rt 和\n'
         '行为转换率 Ch，在散点图上\n'
         '定位教学类型：练习型、讲\n'
         '授型、对话型或混合型，为\n'
         '教学模式诊断提供依据。')

add_card(slide, Inches(9.1), Inches(1.6), Inches(3.8), Inches(2.5),
         '📉 课堂节奏分析',
         '按分钟统计教学活动密度，\n'
         '用折线图展示课堂节奏变化，\n'
         '识别课堂的高潮期和低谷期，\n'
         '帮助教师优化教学环节的时\n'
         '间分配。')

add_card(slide, Inches(0.5), Inches(4.5), Inches(3.8), Inches(2.5),
         '🔥 互动热度分析',
         '按 5 分钟时间段用堆叠柱状\n'
         '图展示各类行为持续时长分\n'
         '布，直观对比不同时间段的\n'
         '教学互动强度变化趋势。')

add_card(slide, Inches(4.8), Inches(4.5), Inches(3.8), Inches(2.5),
         '❓ 提问类型分布',
         '用极坐标图展示不同类型提\n'
         '问的占比，分析教师提问策\n'
         '略的多样性，评估课堂提问\n'
         '的质量和层次分布。')

add_card(slide, Inches(9.1), Inches(4.5), Inches(3.8), Inches(2.5),
         '🔄 教学环节分析',
         '用多线图展示各教学阶段中\n'
         '不同行为类型的占比变化，\n'
         '追踪导入、新授、练习、总\n'
         '结等环节的教学行为演变。')

add_page_number(slide, 7, TOTAL)

# =====================================================
# 第8页：技术难点与解决方案
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT)
add_top_bar(slide, '技术难点与解决方案')

problems = [
    ('Whisper 片段碎片化\n导致大量短片段', 
     '实现智能片段合并算法：将相邻且间隔≤3秒的片段合并，\n'
     '合并后按总时长分配类型，过短片段（<8秒）直接丢弃，\n'
     '有效减少碎片化标注，提升标注结果的合理性。'),
    ('关键词匹配覆盖率低\n大量片段无法标注', 
     '持续扩展关键词库至80+模式，覆盖12个教学行为分类，\n'
     '涵盖真实课堂的主要场景；合并后的未匹配长片段自动\n'
     '标注为"教师讲解"，确保每个有效片段都有标注。'),
    ('Chart.js 画布复用冲突\n同一Canvas重复创建图表', 
     '在 destroyChart 方法中清理 Canvas 2D 上下文和 chartjs\n'
     '内部属性；所有 render 方法用 try-catch 包裹，单个图表\n'
     '出错不影响其他图表展示，保证页面稳定性。'),
    ('处理完成后按钮状态\n未更新', 
     '将 UI 更新逻辑（按钮文字、课程信息表单）移到 setTags\n'
     '触发图表刷新之前执行，避免图表渲染异常中断后续代码，\n'
     '确保按钮状态始终正确更新。'),
]

for i, (problem, solution) in enumerate(problems):
    y = Inches(1.6 + i * 1.45)
    add_shape(slide, Inches(0.5), y, Inches(3.5), Inches(1.2), WHITE)
    add_shape(slide, Inches(0.5), y, Inches(0.08), Inches(1.2), RED)
    add_text_box(slide, Inches(0.8), y + Inches(0.1), Inches(3), Inches(1),
                 problem, font_size=14, color=DARK, bold=True)
    add_shape(slide, Inches(4.2), y, Inches(8.6), Inches(1.2), WHITE)
    add_shape(slide, Inches(4.2), y, Inches(0.08), Inches(1.2), GREEN)
    add_text_box(slide, Inches(4.5), y + Inches(0.1), Inches(8.1), Inches(1),
                 solution, font_size=12, color=DARK)

add_page_number(slide, 8, TOTAL)

# =====================================================
# 第9页：开发过程回顾
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT)
add_top_bar(slide, '开发过程回顾（16天迭代）')

phases = [
    ('第1-2天', '基础搭建', '项目初始化、技术选型、页面布局、导航系统', PRIMARY),
    ('第3-5天', '上传与流水线', '视频上传、FFmpeg处理、Whisper集成、进度轮询', ACCENT),
    ('第6-7天', 'AI 标注引擎', '关键词匹配算法、行为分类体系、标注数据生成', GREEN),
    ('第8-9天', '交互与可视化', '标注编辑管理、Canvas时间轴、视频同步播放', ORANGE),
    ('第10-11天', '多维分析', '弗兰德斯、S-T分析、6种图表、数据聚合计算', RGBColor(0x8B, 0x5C, 0xF6)),
    ('第12-13天', '算法优化', '关键词扩展、片段合并、默认标注策略调整', RED),
    ('第14-15天', 'Bug修复与报告', 'Canvas冲突修复、按钮状态修复、报告页面', RGBColor(0xEC, 0x48, 0x99)),
    ('第16天', '联调交付', '全流程测试、边缘情况修复、文档整理、项目收尾', DARK),
]

for i, (day, title, desc, color) in enumerate(phases):
    row = i // 4
    col = i % 4
    x = Inches(0.5 + col * 3.15)
    y = Inches(1.6 + row * 2.85)
    add_shape(slide, x, y, Inches(2.95), Inches(2.5), WHITE)
    add_shape(slide, x, y, Inches(2.95), Inches(0.55), color)
    add_text_box(slide, x, y + Inches(0.05), Inches(2.95), Inches(0.5),
                 f'{day}：{title}', font_size=15, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.75), Inches(2.65), Inches(1.5),
                 desc, font_size=12, color=DARK)

add_page_number(slide, 9, TOTAL)

# =====================================================
# 第10页：系统功能总览
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT)
add_top_bar(slide, '系统功能模块总览')

modules = [
    ('📤', '上传模块', '拖拽/点击上传\n格式校验\n进度轮询\n流水线触发', PRIMARY),
    ('✏️', '标注模块', 'AI自动标注\n手动编辑标注\n视频同步播放\n标注列表管理', ACCENT),
    ('⏱️', '时间轴模块', 'Canvas绘制\n彩色分层渲染\n缩放与拖拽\n双向视频同步', GREEN),
    ('📊', '分析模块', '弗兰德斯分析\nS-T教学分析\n6种可视化图表\n多维数据聚合', ORANGE),
    ('📋', '报告模块', '结果汇总展示\n教学诊断建议\n打印导出功能\n数据实时刷新', RGBColor(0x8B, 0x5C, 0xF6)),
]

for i, (icon, name, desc, color) in enumerate(modules):
    x = Inches(0.3 + i * 2.55)
    add_shape(slide, x, Inches(1.6), Inches(2.35), Inches(4.5), WHITE)
    add_shape(slide, x, Inches(1.6), Inches(2.35), Inches(0.65), color)
    add_text_box(slide, x, Inches(1.7), Inches(2.35), Inches(0.5),
                 f'{icon} {name}', font_size=18, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(1.95), Inches(3.3),
                 desc, font_size=14, color=DARK, alignment=PP_ALIGN.CENTER)

# 底部：技术栈
add_shape(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8), WHITE)
add_text_box(slide, Inches(0.8), Inches(6.35), Inches(12), Inches(0.7),
             '🛠️ 技术栈：Vite + JavaScript (前端)  |  Node.js + Express (后端)  |  FFmpeg (视频处理)  |  '
             'Whisper.cpp (语音识别)  |  Chart.js (可视化)  |  Canvas API (时间轴)  |  Python-pptx (报告导出)',
             font_size=12, color=DARK, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 10, TOTAL)

# =====================================================
# 第11页：项目成果与亮点
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT)
add_top_bar(slide, '项目成果与亮点')

highlights = [
    ('🎯', '全流程自动化',
     '从视频上传到分析报告，全程无需人工干预。用户上传视频后，系统自动完成音频提取、语音识别、\n'
     '教学行为标注、多维度数据分析，数分钟内生成完整的课堂分析报告。'),
    ('🧠', 'AI 驱动的智能标注',
     '基于 Whisper 语音识别模型和自定义教学行为关键词库（12分类80+模式），结合智能片段合并\n'
     '算法，实现高覆盖率的教学行为自动标注，有效替代人工标注工作。'),
    ('📊', '多维度分析体系',
     '融合弗兰德斯互动分析、S-T教学分析、课堂节奏分析、互动热度分析、提问类型分布、教学环节\n'
     '分析等6个维度，从不同角度全面诊断课堂教学质量。'),
    ('🔄', '实时数据同步架构',
     '基于发布-订阅模式的 DataStore 状态管理，标注数据变更后自动同步到时间轴、分析图表、\n'
     '报告页面等所有模块，保证数据一致性和实时性。'),
]

for i, (icon, title, desc) in enumerate(highlights):
    y = Inches(1.6 + i * 1.4)
    add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(1.15), WHITE)
    add_shape(slide, Inches(0.5), y, Inches(0.08), Inches(1.15), PRIMARY)
    add_text_box(slide, Inches(0.8), y + Inches(0.1), Inches(2), Inches(0.5),
                 f'{icon} {title}', font_size=18, color=PRIMARY, bold=True)
    add_text_box(slide, Inches(0.8), y + Inches(0.55), Inches(11.7), Inches(0.5),
                 desc, font_size=12, color=DARK)

add_page_number(slide, 11, TOTAL)

# =====================================================
# 第12页：总结与展望
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), PRIMARY)
add_shape(slide, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), PRIMARY)

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
             '总结与展望', font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(1.7), Inches(2.3), Inches(0.03), ACCENT)

# 总结
add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(0.5),
             '✅ 项目总结', font_size=22, color=ACCENT, bold=True)

summary_items = [
    '成功构建了从视频上传到智能分析报告的完整工具链',
    '掌握了前端 SPA 架构、后端 API 设计、音视频处理等全栈技能',
    '深入理解了 Whisper 语音识别、关键词匹配、数据聚合等 AI 技术',
    '积累了 Chart.js 数据可视化、Canvas 绘图等前端进阶经验',
    '通过 16 天迭代，经历了从 0 到 1 的完整项目开发流程',
]
for i, item in enumerate(summary_items):
    add_text_box(slide, Inches(2), Inches(2.9 + i * 0.45), Inches(9.5), Inches(0.4),
                 f'▸ {item}', font_size=14, color=LIGHT)

# 展望
add_text_box(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(0.5),
             '🚀 未来展望', font_size=22, color=ACCENT, bold=True)

future_items = [
    '接入更精准的大语言模型（如通义千问、DeepSeek）提升标注准确率',
    '增加多人多课对比分析功能，支持教师群体教学诊断',
    '实现标注数据的导出和与其他教学平台的对接',
    '优化前端性能，支持更长视频（2小时以上）的处理',
]
for i, item in enumerate(future_items):
    add_text_box(slide, Inches(2), Inches(6.0 + i * 0.4), Inches(9.5), Inches(0.35),
                 f'▸ {item}', font_size=13, color=LIGHT)

add_text_box(slide, Inches(1), Inches(7.6), Inches(11), Inches(0.4),
             '感谢聆听  ·  欢迎提问', font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 12, TOTAL)

# ========== 保存 ==========
output_path = 'c:/Users/Administrator/Desktop/项目实践/项目成果汇报.pptx'
prs.save(output_path)
print(f'PPT 已生成：{output_path}')
print(f'共 {len(prs.slides)} 页')
